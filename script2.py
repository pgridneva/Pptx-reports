#!/usr/bin/env python
# coding: utf-8
# number of rows in pass/fail overview slides
num_rows_0 = 15
# rows in table with measurement names and its parameters (2nd on slide)
num_rows_2 = 12
# rows in table "major contributor to variations"
num_rows_3 = 4
# rows in the last slides (used analysis data from csv)
num_rows_end = 20
num_templates = 10

lsl_nan = -7874016000000000000.0 # specific value for nan-value of spes in csv
usl_nan = 7874016000000000000.0

tabl3_if_string = 1 # if in cells = 0 

# position of histogram-pic on slide, left/top - x/y coordinates of top-left corner of pic 
# then setting width to adjust the size of pic
from pptx.util import Inches, Pt
left = Inches(9.5)
top = Inches(5)
width = Inches(3)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.creator import PPTXCreator
from pptx_tools.position import PPTXPosition
from pptx_tools.table_style import PPTXTableStyle
from pptx_tools.paragraph_style import PPTXParagraphStyle
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import pandas as pd
import six
import copy
import numpy as np
import os

# formatting font in tables, here one can change bold/italic/size/font in parenthesis
# two styles for different font size

table_style = PPTXTableStyle()
table_style.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=12)

table_style1 = PPTXTableStyle()
table_style1.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=10)

# to middle-align text frames
paragraph_style = PPTXParagraphStyle()
paragraph_style.set(alignment=PP_PARAGRAPH_ALIGNMENT.CENTER)

# function to copy slide with given index and place it to the end of the presentation
def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[5]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )
    return copied_slide

# data to be put in table (x-measurements are dropped)

# used data, Book1
df0 = pd.read_excel('Book1.xlsx', sheet_name='Measures', header=0)
# dropping duplicated lines (appear because of xlsx formatting)
df0 = df0.drop_duplicates(subset=['PART', 'Measure'])
df0.reset_index(inplace=True, drop=True)

# deleting inactive records and x-measurements
wrong_ind = []
for i in range(len(df0)):
    if 'x' in df0.iloc[i,1]:
        wrong_ind.append(i)
    if df0.loc[i, 'As Output'] == 0:
        wrong_ind.append(i)
    if df0.loc[i, 'Active'] == 0:
        wrong_ind.append(i)
df0.drop(wrong_ind, inplace=True)
df0.reset_index(inplace=True, drop=True)

# deleting records where description is empty
df0.dropna(subset=['Description'], inplace=True)
df0.reset_index(inplace=True, drop=True)

# leaving only necessary columns
df0 = df0[['Description', 'Type', 'Measure', 'PART', 'USpecL', 'LSpecL', 'Relative/Abs']]

# checking if where is a description with DRA number in the beginning
if df0.loc[0, 'Description'][0] == '#':
    with_index = True
else:
    with_index = False
    
# adding column with DRA-num
if with_index:
    for i in range(len(df0)):
        desc = str(df0.loc[i, 'Description']).split(sep=':')
        df0.loc[i,'idx'] = int(desc[0][1:])
        
# reading csv
df1 = pd.read_excel('tabl2.xlsx', header=12)
# not too output -0.0000 in Nominal further
for i in range(len(df1)):
    if abs(df1.loc[i, 'Nominal']) < 0.0001:
        df1.loc[i, 'Nominal'] = 0
 
# keeping all columns from csv to put data on the last slides       
df1_init = df1[['Name',	'Nominal',	'Mean',	'6-Sigma', 'Min', 'Max'	,'Range','LSL','USL',
               'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Type', 'Est.Low', 'Est.High', 'Est.Range']]
# total number of records
num_meas_total = len(df1_init)

# dropping x-measurements from csv
wrong_ind = []
for i in range(len(df1)):
    if 'x' in df1.iloc[i,1]:
        wrong_ind.append(i)
df1.drop(wrong_ind, inplace=True)
df1.reset_index(inplace=True, drop=True)
df1 = df1[['Name', 'Description', 'Nominal', 'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range']]
# УБРАТЬ РЕШЕТКИ В ДВУХ СТРОКАХ НИЖЕ, ЕСЛИ ИСПОЛЬЗУЕТСЯ ALIAS В ИЗМЕРЕНИЯХ
#for i in range(len(df1)):
#    df1.loc[i, 'Name'] = df1.loc[i, 'Name'][3:]

# merging data about certain measurements
df_01 = df0.merge(df1, left_on='Measure', right_on='Name')
df_01 = df_01.drop(['Description_x', 'Measure'], axis=1)

# if DRA-num is in description, sort all data by it, also adding Est.Low.C and Est.High.C columns
# for convenience, order of columns is changed
if with_index:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs', 'idx']]
    df_01['Est.Low.C'] = df_01['Est.Low'] - df_01['Nominal']
    df_01['Est.High.C'] = df_01['Est.High'] - df_01['Nominal']
    df_01 = df_01.sort_values(by='idx', ignore_index=True)
else:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs']]
    df_01['Est.Low.C'] = df_01['Est.Low'] - df_01['Nominal']
    df_01['Est.High.C'] = df_01['Est.High'] - df_01['Nominal']

# reading rss omitting lines in the beginning
df2 = pd.read_excel('tabl3.xlsx', skiprows=10)

# array of tables corresponding to measurements
# and array of descriptions for convenience
meas_list = []
meas_names = []

for name in df_01['Description_y'].unique():
    meas_names.append(name)
    df_i = df_01[df_01['Description_y'] == name]
    df_i.reset_index(inplace=True, drop=True)
    df_i = df_i.sort_values(by='Name', ignore_index=True)
    meas_list.append(df_i)

# below the code to decide about worst meaasurements

# list of worst measurements (indexes correspond to meas_names)
contr_choice = []
# lists of fallouts of worst measurements (indexes correspond to meas_names)
fallouts = []

# first look for meas. with max tot-out, if it is not unique then choosing the meas. with max est.range
for df in meas_list:
    if len(df) == 1:
        point = df.loc[0, 'Name']
    else:
        ind = 0 # how many records have value of max tot-out
        max_t_o = df['Tot-OUT'].max()
        for val in df['Tot-OUT']:
            if val == max_t_o:
                ind += 1
        if ind > 1:
            point = df.loc[df['Est.Range'].idxmax(), 'Name']
        else:
            point = df.loc[df['Tot-OUT'].idxmax(), 'Name']
            
    contr_choice.append(point)
    fallouts.append(df['Tot-OUT'].max())

# working with rss
# finding empty lines 
# getting array of indexes where records start
newrec_arr = []
for i in range(len(df2)):
    if str(df2.iloc[i, 0]) == 'nan':
        newrec_arr.append(i)

l = newrec_arr
l_mod = [0] + l + [len(df2) + 1]  

# array of dfs on every measurement name
list_dfs = [df2.iloc[l_mod[n]:l_mod[n+1]] for n in range(len(l_mod)-1)]

# deleting first (empty) line in records
for i in range(1, len(list_dfs)):
    list_dfs[i].reset_index(inplace=True, drop=True)
    list_dfs[i].drop([0], inplace=True)
    list_dfs[i].reset_index(inplace=True, drop=True)

# code for creating dfs with contributors-info on every measurement
contr_names_list = [] # contributors
sublist_dfs = [] # array of such dfs

for df in list_dfs:
    # в случае ALIAS используем следующую строку, комментируем через одну
    # contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][3:-1])
    contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][:-1])
    num_recs = int((len(df) - 5) / 2)
    
    sub_df = pd.DataFrame(index=range(num_recs), columns=['contr_name', 'type', 'tol_range', 'contr'])
    for i in range(num_recs):
        if tabl3_if_string == 1:
             str_i = df.iloc[4 + i * 2, 0].split() # 4 initial lines, then 2 lines per 1 contributor 
             # big sring cell is divided to list of strings with white space as delimiter
             
             sub_df.loc[i, 'tol_range'] = str_i[4][2:] # range
             sub_df.loc[i, 'contr'] = str_i[5] # contribution
            
             typ = str_i[1][:2]
             if typ == 'CR':
                 sub_df.loc[i, 'type'] = 'Positional'
             elif typ == 'CH' or typ == 'CP':
                 sub_df.loc[i, 'type'] = 'Size'
             else:
                 sub_df.loc[i, 'type'] = 'Linear'
    
             if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan':
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
             else:
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]
                 
        else: # we just take necessary values from cells in row
            str_i = df.iloc[4 + i * 2, 1]
            sub_df.loc[i, 'tol_range'] = df.iloc[4 + i * 2, 4][2:] # range
            sub_df.loc[i, 'contr'] = str(df.iloc[4 + i * 2, 5]) # contribution
            
            typ = str_i[:1]
            if typ == 'CR':
                sub_df.loc[i, 'type'] = 'Positional'
            elif typ == 'CH' or typ == 'CP':
                sub_df.loc[i, 'type'] = 'Size'
            else:
                sub_df.loc[i, 'type'] = 'Linear'
    
            if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan': # we stick together 2 cells here if contr's name is not empty 
                sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
            else:
                sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]
    sublist_dfs.append(sub_df)   

# opening existing presentation (template with tables)

pres = Presentation('boe_template.pptx')
pres.slide_width = Inches(13.333333) # widescreen
pres.slide_height = Inches(7.5)

slides = [slide for slide in pres.slides] # array of slide-objects

# reading template_choice list
num_req = len(meas_list) # number of req = number of slides
links = pd.read_excel('template_choice.xlsx', header=0)

# calculating number of slides for pass/fail overview
if num_req % num_rows_0 == 0:
    init_slides_num = int(num_req / num_rows_0)
else:
    init_slides_num = num_req // num_rows_0 + 1

# filling slides with descriptions list   
for i in range(init_slides_num):
    copied_slide = duplicate_slide(pres, 0)  
    slide_i = pres.slides[num_templates + i]
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]
    for k in range(num_rows_0):
        if k <= num_req - num_rows_0 * i - 1:
            if with_index:
                tables_i[0].table.cell(k + 1, 0).text = str(meas_names[i * num_rows_0 + k].split(':')[0])
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_0 + k].split(':')[1][1:])
            else:
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_0 + k])
            tables_i[0].table.cell(k + 1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
            tables_i[0].table.cell(k + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER          
            tables_i[0].table.cell(k + 1, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            lsl, usl = '{:.4f}'.format(meas_list[i * num_rows_0 + k].iloc[0,5]), '{:.4f}'.format(meas_list[i * num_rows_0 + k].iloc[0,4])
            tables_i[0].table.cell(k + 1, 2).text = str(lsl) + '/' + str(usl) 
            tables_i[0].table.cell(k + 1, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
            fallout = fallouts[i * num_rows_0 + k]
            
            if str(lsl) == 'nan':
                tables_i[0].table.cell(k + 1, 3).text = 'N/A'
                tables_i[0].table.cell(k + 1, 2).text = 'N/A'
            else:
                tables_i[0].table.cell(k + 1, 3).text = '{:.1%}'.format(fallout)   
                fill = tables_i[0].table.cell(k + 1, 3).fill
                fore_color = fill.solid()
                if fallout == 0:
                    tables_i[0].table.cell(k + 1, 3).text = 'Pass'
                    fill.fore_color.rgb = RGBColor(198, 239, 206)
                elif fallout < 0.05:
                    fill.fore_color.rgb = RGBColor(255, 235, 156)
                elif fallout >= 0.05:
                    fill.fore_color.rgb = RGBColor(255, 199, 206)
            tables_i[0].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[0].table.cell(k+1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER    
            tables_i[0].table.cell(k+1, 3).vertical_anchor = MSO_ANCHOR.MIDDLE                                   
    table_style.write_shape(tables_i[0])

dop_sl = 0
# 1 description - 1 slide plus several more if list of measurements doesn't fit to 1 slide
# dop_sl is the sum of added slides to output all the measurements in list

# slides with measurements under 1 description and worst measurement
for i in range(num_req):
    num_meas_i = len(meas_list[i])
    if num_meas_i <= num_rows_2: # all measurements are being placed on 1 slide
        template_type = links['template'][i]
        copied_slide = duplicate_slide(pres, template_type)
        slide_loc = num_templates + init_slides_num + i + dop_sl
        #print(i, slide_loc)
        slide_i = pres.slides[slide_loc] 
        tables_i = [shape for shape in slide_i.shapes if shape.has_table]
        
        if with_index:
            req = meas_list[i].iloc[0,2].split(':')[0]+' -'+meas_list[i].iloc[0,2].split(':')[1]
        else:
            req = meas_list[i].iloc[0,2]
        mes_type = meas_list[i].iloc[0, 3]
        tables_i[0].table.cell(0,2).text = str(req)
        tables_i[0].table.cell(1,2).text = str(mes_type)
        if meas_list[i].loc[0, 'Relative/Abs'] == 0:
            tables_i[0].table.cell(1,8).text = 'Absolute' 
        if meas_list[i].loc[0, 'Relative/Abs'] == 1:
            tables_i[0].table.cell(1,8).text = 'Relative'
        lsl, usl = '{:.4f}'.format(meas_list[i].iloc[0,5]), '{:.4f}'.format(meas_list[i].iloc[0,4])
        tables_i[0].table.cell(1,6).text = str(lsl) + '/' + str(usl)   
        if str(lsl) == 'nan':
            tables_i[0].table.cell(1, 6).text = 'N/A'
        tables_i[0].table.cell(1,6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # just filling table according to number of template
        if links['template'][i] == 1:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Nominal'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 2:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 3:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 4:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 5:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 6:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])       
        elif links['template'][i] == 7:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Nominal'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 8:   
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Nominal'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])

                
        contr = contr_choice[i]
        loc = 0
        for m in range(len(contr_names_list)):
            if contr_names_list[m] == contr:
                loc = m
        subtabl = sublist_dfs[loc] 
        tables_i[2].table.cell(0, 0).text = 'Major Contributors to Variations ({})'.format(contr)
        len_tabl = len(subtabl) if len(subtabl) < num_rows_3 else num_rows_3
        for k in range(len_tabl):
            tables_i[2].table.cell(k+1, 0).text = str(subtabl.iloc[k, 0])
            tables_i[2].table.cell(k+1, 1).text = str(subtabl.iloc[k, 1])
            tables_i[2].table.cell(k+1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[2].table.cell(k+1, 2).text = '{:.4f}'.format(float(subtabl.iloc[k, 2]))
            tables_i[2].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if tabl3_if_string == 0:
                tables_i[2].table.cell(k+1, 3).text = '{:.1%} '.format(float(subtabl.iloc[k, 3]))
            else:
                tables_i[2].table.cell(k+1, 3).text = '{}% '.format(subtabl.iloc[k, 3][:-4])
            tables_i[2].table.cell(k+1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        table_style1.write_shape(tables_i[0])
        table_style1.write_shape(tables_i[1])
        paragraph_style.write_shape(tables_i[1])
        for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT        
        table_style1.write_shape(tables_i[2])

        for filename in os.listdir('images/'):
            if contr in filename:
                img_path = os.path.join('images/', filename)        
        img_i = slide_i.shapes.add_picture(img_path, left, top, width)
        
    else: # if measurements need to be placed at several slides
        if num_meas_i % num_rows_2 == 0:
            slides_per_req = int(num_meas_i / num_rows_2)
        else:
            slides_per_req = num_meas_i // num_rows_2 + 1
            
        for j in range(slides_per_req):
            template_type = links['template'][i]
            copied_slide = duplicate_slide(pres, template_type)  
            slide_loc = num_templates + init_slides_num + i + dop_sl + j
            slide_i = pres.slides[slide_loc] 
            #print(i, slide_loc)
            tables_i = [shape for shape in slide_i.shapes if shape.has_table]
            
            if with_index:
                 req = meas_list[i].iloc[0,2].split(':')[0]+' -'+meas_list[i].iloc[0,2].split(':')[1]
            else:
                req = meas_list[i].iloc[0,2]
            mes_type = meas_list[i].iloc[0, 3]
            tables_i[0].table.cell(0,2).text = str(req)
            tables_i[0].table.cell(1,2).text = str(mes_type)
            if meas_list[i].loc[0, 'Relative/Abs'] == 0:
                tables_i[0].table.cell(1,8).text = 'Absolute' 
            if meas_list[i].loc[0, 'Relative/Abs'] == 1:
                tables_i[0].table.cell(1,8).text = 'Relative'
            lsl, usl = '{:.4f}'.format(meas_list[i].iloc[0,5]), '{:.4f}'.format(meas_list[i].iloc[0,4])
            if str(lsl) == 'nan':
                tables_i[0].table.cell(1, 6).text = 'N/A'
            else:
                tables_i[0].table.cell(1,6).text = str(lsl) + '/' + str(usl) 
            tables_i[0].table.cell(1,6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            contr = contr_choice[i]
            loc = 0
            for m in range(len(contr_names_list)):
                if contr_names_list[m] == contr:
                    loc = m
            subtabl = sublist_dfs[loc] 
            tables_i[2].table.cell(0, 0).text = 'Major Contributors to Variations ({})'.format(contr)
            len_tabl = len(subtabl) if len(subtabl) < num_rows_3 else num_rows_3
            for k in range(len_tabl):
                tables_i[2].table.cell(k+1, 0).text = str(subtabl.iloc[k, 0])
                #tables_i[2].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                tables_i[2].table.cell(k+1, 1).text = str(subtabl.iloc[k, 1])
                tables_i[2].table.cell(k+1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                tables_i[2].table.cell(k+1, 2).text = '{:.4f}'.format(float(subtabl.iloc[k, 2]))
                tables_i[2].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                if tabl3_if_string == 0:
                    tables_i[2].table.cell(k+1, 3).text = '{:.1%} '.format(float(subtabl.iloc[k, 3]))
                else:
                    tables_i[2].table.cell(k+1, 3).text = '{}% '.format(subtabl.iloc[k, 3][:-4])
                tables_i[2].table.cell(k+1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            for filename in os.listdir('images/'):
                if contr in filename:
                    img_path = os.path.join('images/', filename)        
            img_i = slide_i.shapes.add_picture(img_path, left, top, width)
            
            if j != slides_per_req - 1:  
                meas_list_j = meas_list[i].iloc[num_rows_2 * j: num_rows_2 * (j+1), ]
            else:
                meas_list_j = meas_list[i].iloc[num_rows_2 * j: , ]
            meas_list_j.reset_index(inplace=True, drop=True)
            
            if links['template'][i] == 1: 
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.iloc[k, 0])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.iloc[k, 6], 4)
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.iloc[k, 7])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.iloc[k, 8])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.iloc[k, 9])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.iloc[k, 10])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.iloc[k, 11])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.iloc[k, 14])
                    tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list_j.iloc[k, 15])
                    tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list_j.iloc[k, 12])
            if links['template'][i] == 2:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 3:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 4:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 5:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 6:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])       
            elif links['template'][i] == 7:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Nominal'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 8:   
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Nominal'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])  
                    

            table_style1.write_shape(tables_i[0])
            table_style1.write_shape(tables_i[1])
            paragraph_style.write_shape(tables_i[1])
            for k in range(len(meas_list_j)):
                tables_i[1].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            table_style1.write_shape(tables_i[2])

   
        dop_sl += slides_per_req - 1    
   
# all slides we've generated so far         
num_slides_main = num_templates + num_req + init_slides_num + dop_sl

# number of slides needed for data from csv
num_end_slides = int(num_meas_total / num_rows_end) if num_meas_total % num_rows_end == 0 else num_meas_total // num_rows_end + 1

# tables from csv being put in the end
for i in range(num_end_slides):
    copied_slide = duplicate_slide(pres, num_templates-1)  
    slide_i = pres.slides[num_slides_main + i]
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]    
    
    # we divide our csv-df to small dfs to be put on slides and fill table with them
    if i != num_end_slides - 1:  
        list_i = df1_init.iloc[num_rows_end * i: num_rows_end * (i+1), ] # last dataframe
    else:
        list_i = df1_init.iloc[num_rows_end * i: , ]
    list_i.reset_index(inplace=True, drop=True)
        
                    
    for k in range(len(list_i)):
        tables_i[0].table.cell(k+1, 0).text = str(list_i.loc[k, 'Name'])
        tables_i[0].table.cell(k+1, 1).text = '{:.4f}'.format(list_i.loc[k, 'Nominal'])
        tables_i[0].table.cell(k+1, 2).text = '{:.4f}'.format(list_i.loc[k, 'Mean'])
        tables_i[0].table.cell(k+1, 3).text = '{:.4f}'.format(list_i.loc[k, '6-Sigma'])
        tables_i[0].table.cell(k+1, 4).text = '{:.4f}'.format(list_i.loc[k, 'Min'])
        tables_i[0].table.cell(k+1, 5).text = '{:.4f}'.format(list_i.loc[k, 'Max'])
        tables_i[0].table.cell(k+1, 6).text = '{:.4f}'.format(list_i.loc[k, 'Range'])
        lsl, usl = list_i.loc[k, 'LSL'], list_i.loc[k, 'USL']
        if lsl == lsl_nan:
            tables_i[0].table.cell(k+1, 7).text = 'N/A'
        else:
            tables_i[0].table.cell(k+1, 7).text = '{:.4f}'.format(lsl) 
        if usl == usl_nan:    
            tables_i[0].table.cell(k+1, 8).text = 'N/A'
        else:
            tables_i[0].table.cell(k+1, 8).text = '{:.4f}'.format(usl)
        tables_i[0].table.cell(k+1, 9).text = '{:.1%}'.format(list_i.loc[k,'L-OUT'])
        tables_i[0].table.cell(k+1, 10).text = '{:.1%}'.format(list_i.loc[k,'H-OUT'])
        tables_i[0].table.cell(k+1, 11).text = '{:.1%}'.format(list_i.loc[k, 'Tot-OUT'])
        etype = list_i.loc[k, 'Est.Type']
        if etype == '   Normal':
            tables_i[0].table.cell(k+1, 12).text = etype[3:]
        else:
            tables_i[0].table.cell(k+1, 12).text = etype
        tables_i[0].table.cell(k+1, 13).text = '{:.4f}'.format(list_i.loc[k, 'Est.Low'])
        tables_i[0].table.cell(k+1, 14).text = '{:.4f}'.format(list_i.loc[k, 'Est.High'])
        tables_i[0].table.cell(k+1, 15).text = '{:.4f}'.format(list_i.loc[k, 'Est.Range'])
        
    paragraph_style.write_shape(tables_i[0])
    # meas name alignment
    for k in range(len(list_i)):
                tables_i[0].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT        
    table_style1.write_shape(tables_i[0])
    
# deleting several template slides in the beginning, better do not change
for i in range(num_templates-1,-1,-1) : 
    rId = pres.slides._sldIdLst[i].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[i]

pres.save('pres_upd.pptx') # saving file   




