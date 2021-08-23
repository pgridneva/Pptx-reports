#!/usr/bin/env python
# coding: utf-8

# number of rows in list of measurements slide (and pass/fail overview)
num_rows_meas = 14 
# number of rows in table with all parameters given (slide 4 in template)
num_rows_big_tabl = 21 
# number of template slide in pptx file, better add new one in the end to avoid script failure
num_templates = 6
# rows in table "major contributor to variations"
num_rows_contr = 5
# rows in measurements descriprion slide (2nd in template)
num_rows_groups = 10
# rows in the last slides (used analysis data from csv)
num_rows_end = 20

lsl_nan = -7874016000000000000.0 # specific value for nan-value of spes in csv
usl_nan = 7874016000000000000.0

tabl3_if_string = 1 # in rss if in cells = 0 

# position of histogram-pic on slide, left/top - x/y coordinates of top-left corner of pic 
# then setting width to adjust the size of pic
from pptx.util import Inches, Pt
left = Inches(0.5) 
top = Inches(5)
width = Inches(3)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.creator import PPTXCreator
from pptx_tools.position import PPTXPosition
from pptx_tools.table_style import PPTXTableStyle
from pptx_tools.paragraph_style import PPTXParagraphStyle
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import pandas as pd
import six
import copy
import numpy as np
import os
from pptx_tools.templates import TemplateExample
pp = PPTXCreator(TemplateExample())

# formatting font in tables, here one can change bold/italic/size/font in parenthesis
# two styles for different font size
table_style_12 = PPTXTableStyle()
table_style_12.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=12)

table_style_10 = PPTXTableStyle()
table_style_10.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=10)

# header on slide 5 (description of worst measurement and its name), blue color, bold
style_contr_header = PPTXTableStyle()
style_contr_header.font_style = PPTXFontStyle().set(bold=True, name="Arial", size=14, color_rgb=(0, 0, 255))

# yellow text frames (name of measurements, contribution procents, specs)
font_yllw_box = PPTXFontStyle()
font_yllw_box.set(size = 12, bold=False, name="Arial")

# to middle-align text frames
paragraph_style = PPTXParagraphStyle()
paragraph_style.set(alignment=PP_PARAGRAPH_ALIGNMENT.CENTER)

# function to copy slide with given index and place it to the end of the presentation
def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[5]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )
    return copied_slide

# used data, Book1
df0 = pd.read_excel('Book1.xlsx', sheet_name='Measures', header=0)

# dropping duplicated lines (appear because of xlsx formatting)
df0 = df0.drop_duplicates(subset=['PART', 'Measure'])
df0.reset_index(inplace=True, drop=True)

# deleting inactive records and x-measurements
wrong_ind = []
for i in range(len(df0)):
    if 'x' in df0.loc[i, 'Measure']:
        wrong_ind.append(i)
    if df0.loc[i, 'As Output'] == 0:
        wrong_ind.append(i)
    if df0.loc[i, 'Active'] == 0:
        wrong_ind.append(i)       
df0.drop(wrong_ind, inplace=True)
df0.reset_index(inplace=True, drop=True)

# deleting records where description is empty
df0.dropna(subset=['Description'], inplace=True)
df0.reset_index(inplace=True, drop=True)

# leaving only necessary columns
df0 = df0[['Description', 'Type', 'Measure', 'PART', 'USpecL', 'LSpecL', 'Relative/Abs']]

# checking if where is a description with DRA number in the beginning
if df0.loc[0, 'Description'][0] == '#':
    with_index = True
else:
    with_index = False

# adding column with DRA-num
if with_index:
    for i in range(len(df0)):
        desc = str(df0.loc[i, 'Description']).split(sep=':')
        df0.loc[i,'idx'] = int(desc[0][1:])

# reading csv
df1 = pd.read_excel('tabl2.xlsx', header=12)
# not too output -0.0000 in Nominal further
for i in range(len(df1)):
    if abs(df1.loc[i, 'Nominal']) < 0.0001:
        df1.loc[i, 'Nominal'] = 0

# keeping all columns from csv to put data on the last slides
df1_init = df1[['Name',	'Nominal',	'Mean',	'6-Sigma', 'Min', 'Max'	,'Range','LSL','USL',
               'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Type', 'Est.Low', 'Est.High', 'Est.Range']]

# total number of records
num_meas_total = len(df1_init)

# dropping x-measurements from csv
wrong_ind = []
for i in range(len(df1)):
    if 'x' in df1.loc[i, 'Name']:
        wrong_ind.append(i)
df1.drop(wrong_ind, inplace=True)
df1.reset_index(inplace=True, drop=True)

# leaving only necessary columns
df1 = df1[['Name', 'Description', 'Nominal', 'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range']]

# if ALIAS, uncomment these 2 lines
#for i in range(len(df1)):
#    df1.loc[i, 'Name'] = df1.loc[i, 'Name'][3:]

# merging data from Book1 and csv about certain measurements 
df_01 = df0.merge(df1, left_on='Measure', right_on='Name')
df_01 = df_01.drop(['Description_x', 'Measure'], axis=1)

# if DRA-num is in description, sort all data by it, also adding Est.Low.C and Est.High.C columns
# for convenience, order of columns is changed
if with_index:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs', 'idx']]
    df_01 = df_01.sort_values(by='idx', ignore_index=True)
else:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs']]
    
df_01['Est.Low.C'] = df_01['Est.Low'] - df_01['Nominal']
df_01['Est.High.C'] = df_01['Est.High'] - df_01['Nominal']

# reading rss omitting lines in the beginning
df2 = pd.read_excel('tabl3.xlsx', skiprows=10)

# creating array of dataframes corresponding to one description
meas_list = []
# list of descriptions
meas_names = []
for name in df_01['Description_y'].unique():
    meas_names.append(name)
    df_i = df_01[df_01['Description_y'] == name]
    df_i.reset_index(inplace=True, drop=True)
    df_i = df_i.sort_values(by='Name', ignore_index=True) # sorting by measuremments name
    meas_list.append(df_i)

# below the code to decide about worst meaasurements

# list of worst measurements (indexes correspond to meas_names)
contr_choice = []

# indexes show the position of worst measuremnt in df of its description  
contr_locs = []

# lists of fallouts and ranges of worst measurements (indexes correspond to meas_names)
fallouts = []
ranges = []

# first look for meas. with max tot-out, if it is not unique then choosing the meas. with max est.range
for df in meas_list:
    if len(df) == 1:
        point = df.loc[0, 'Name']
        loc = 0
    else:
        ind = 0 # how many records have value of max tot-out
        max_t_o = df['Tot-OUT'].max()
        for val in df['Tot-OUT']:
            if val == max_t_o:
                ind += 1
        if ind > 1:
            loc = df['Est.Range'].idxmax()
            point = df.loc[loc, 'Name']
        else:
            loc = df['Tot-OUT'].idxmax()
            point = df.loc[loc, 'Name']   
    contr_locs.append(loc)        
    contr_choice.append(point)
    fallouts.append(df['Tot-OUT'].max())
    ranges.append(df['Est.Range'].max())

# working with rss
    
# finding empty lines which divide records   
# getting array of indexes where records start
newrec_arr = []
for i in range(len(df2)):
    if str(df2.iloc[i, 0]) == 'nan':
        newrec_arr.append(i)
l = newrec_arr
l_mod = [0] + l + [len(df2) + 1]  

# array of dfs on every measurement name
list_dfs = [df2.iloc[l_mod[n]:l_mod[n+1]] for n in range(len(l_mod)-1)]

# deleting first (empty) line in records
for i in range(1, len(list_dfs)):
    list_dfs[i].reset_index(inplace=True, drop=True)
    list_dfs[i].drop([0], inplace=True)
    list_dfs[i].reset_index(inplace=True, drop=True)

# code for creating dfs with contributors-info on every measurement
contr_names_list = [] # contributors
sublist_dfs = [] # array of such dfs

for df in list_dfs:
    # in case of ALIAS uncomment this:
    # contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][3:-1])
    contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][:-1])
    num_recs = int((len(df) - 5) / 2) # num of contributors in meas-record
    
    sub_df = pd.DataFrame(index=range(num_recs), columns=['contr_name', 'type', 'tol_range', 'contr'])
    for i in range(num_recs):
        if tabl3_if_string == 1:
             str_i = df.iloc[4 + i * 2, 0].split() # 4 initial lines, then 2 lines per 1 contributor 
             # big sring cell is divided to list of strings with white space as delimiter
             sub_df.loc[i, 'tol_range'] = str_i[4][2:] # range
             sub_df.loc[i, 'contr'] = str_i[5] # contribution
            
             typ = str_i[1][:2]
             if typ == 'CR':
                 sub_df.loc[i, 'type'] = 'Positional'
             elif typ == 'CH' or typ == 'CP':
                 sub_df.loc[i, 'type'] = 'Size'
             else:
                 sub_df.loc[i, 'type'] = 'Linear'
    
             if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan':
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
             else:
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]
                 
        else: # we just take necessary values from cells in row
            str_i = df.iloc[4 + i * 2, 1]
            sub_df.loc[i, 'tol_range'] = df.iloc[4 + i * 2, 4][2:] # range
            sub_df.loc[i, 'contr'] = str(df.iloc[4 + i * 2, 5]) # contribution
            
            typ = str_i[:1]
            if typ == 'CR':
                sub_df.loc[i, 'type'] = 'Positional'
            elif typ == 'CH' or typ == 'CP':
                sub_df.loc[i, 'type'] = 'Size'
            else:
                sub_df.loc[i, 'type'] = 'Linear'
    
            if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan': # we stick together 2 cells here if contr's name is not empty 
                sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
            else:
                sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]
    sublist_dfs.append(sub_df)   

# opeining pptx template
pres = Presentation('template_latch.pptx')
slides = [slide for slide in pres.slides] # array of slide-objects
num_req = len(meas_list) # number of descriptions

# calculating number of slides "meas-s list" and "pass/fail overview"
if num_req % num_rows_meas == 0:
    init_slides_num = int(num_req / num_rows_meas)
else:
    init_slides_num = num_req // num_rows_meas + 1
 
# filling slides with descriptions list    
for i in range(init_slides_num):
    copied_slide = duplicate_slide(pres, 0)  # choosing right template
    slide_i = pres.slides[num_templates + i]
    # table-type objects on the slide
    tables_i = [shape for shape in slide_i.shapes if shape.has_table] 
    for k in range(num_rows_meas):
        if k <= num_req - num_rows_meas * i - 1:        
            if with_index:
                # filling cells with DRA-num and description
                tables_i[0].table.cell(k + 1, 0).text = str(meas_names[i * num_rows_meas + k].split(sep=':')[0])
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_meas + k].split(sep=':')[1][1:])
            else:
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_meas + k])
            # formatting text alignment, horizontal, middle
            tables_i[0].table.cell(k + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # vertical alignment (center of the cell)
            tables_i[0].table.cell(k + 1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE               
            tables_i[0].table.cell(k + 1, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            
            lsl, usl = '{:.4f}'.format(meas_list[i * num_rows_meas + k].loc[0,'LSpecL']), '{:.4f}'.format(meas_list[i * num_rows_meas + k].loc[0,'USpecL'])
            tables_i[0].table.cell(k + 1, 2).text = str(lsl) + '/' + str(usl) 
            tables_i[0].table.cell(k + 1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[0].table.cell(k + 1, 2).vertical_anchor = MSO_ANCHOR.MIDDLE             
            # just writing to Specs cell N/A instead of nan
            if str(lsl) == 'nan':
                tables_i[0].table.cell(k + 1, 2).text = 'N/A'
    # setting font of text in table
    table_style_12.write_shape(tables_i[0])          
    
# if DRA-num is in description, making slides with DRA-grouped descriptions
if with_index:
    values = df_01['idx'].unique() # list of unique DRA-numbers
    num_idx = len(values) # how many such slides we'll need    
    for i in range(num_idx):
        copied_slide = duplicate_slide(pres, 1) # !!! here index of slide in template starts from 0
        slide_i = pres.slides[num_templates + init_slides_num + i]
        tables_i = [shape for shape in slide_i.shapes if shape.has_table]  
        # choosing certain DRA-num
        val = values[i]
        k = 0 # how many descriptions in this group
        for j in range(len(meas_list)): # checking all descriptions 
            if meas_list[j].loc[0, 'idx'] == val:    
                # filling table row by row, first DRA-num, then description
                tables_i[0].table.cell(k+1, 0).text = str(meas_list[j].loc[0,'Description_y'] .split(sep=':')[0])
                tables_i[0].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                tables_i[0].table.cell(k+1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE   
                tables_i[0].table.cell(k+1, 1).text = str(meas_list[j].loc[0,'Description_y'].split(sep=':')[1][1:])
                tables_i[0].table.cell(k+1, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
                # Specs
                lsl, usl = '{:.4f}'.format(meas_list[j].loc[0,'LSpecL']), '{:.4f}'.format(meas_list[j].loc[0,'USpecL'])
                tables_i[0].table.cell(k+1, 2).text = str(lsl) + '/' + str(usl) 
                if str(lsl) == 'nan':
                    tables_i[0].table.cell(k+1, 2).text = 'N/A'
                tables_i[0].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                tables_i[0].table.cell(k+1, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
                k += 1
                
        # merging cells with common index
        if k > 1:
            tables_i[0].table.cell(1, 0).merge(tables_i[0].table.cell(k, 0))
            tables_i[0].table.cell(1, 0).text = '#' + str(int(val))
            tables_i[0].table.cell(1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
        table_style_12.write_shape(tables_i[0]) 
        
        # text-frame objects in slide
        title_shapes = [shape for shape in slide_i.shapes if shape.has_text_frame]
        titles = [shape for shape in title_shapes if shape.has_text_frame]
        # filling them with "Req i"
        for i in range(5):
            # starting from 3 as there's also text shapes at the bottom of the slide
            titles[3 + i].text = 'Req ' + str(int(val))
            font_yllw_box.write_paragraph(titles[3 + i].text_frame.paragraphs[0])
else: #  if no DRA-num, we add slide-measurement
    num_idx = num_req
    for i in range(num_idx):
        copied_slide = duplicate_slide(pres, 1) # !!! here index of slide in template starts from 0
        slide_i = pres.slides[num_templates + init_slides_num + i]
        tables_i = [shape for shape in slide_i.shapes if shape.has_table] 
        tables_i[0].table.cell(1, 1).text = str(meas_list[i].loc[0,'Description_y'])
        tables_i[0].table.cell(1, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
        # Specs
        lsl, usl = '{:.4f}'.format(meas_list[i].loc[0,'LSpecL']), '{:.4f}'.format(meas_list[i].loc[0,'USpecL'])
        tables_i[0].table.cell(1, 2).text = str(lsl) + '/' + str(usl) 
        if str(lsl) == 'nan':
            tables_i[0].table.cell(1, 2).text = 'N/A'
        tables_i[0].table.cell(1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tables_i[0].table.cell(1, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
        table_style_12.write_shape(tables_i[0]) 

        
# slides for pass/fail overview, quite similar to measurements list
for i in range(init_slides_num):
    copied_slide = duplicate_slide(pres, 2) # choosing necessary slide from template
    slide_i = pres.slides[num_templates + init_slides_num + num_idx + i]
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]
    for k in range(num_rows_meas):
        if k <= num_req - num_rows_meas * i - 1: # here we can stop filling rows below
            if with_index:
                tables_i[0].table.cell(k + 1, 0).text = str(meas_names[i * num_rows_meas + k].split(sep=':')[0])
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_meas + k].split(sep=':')[1][1:])
            else:
                tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_meas + k])
            tables_i[0].table.cell(k + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[0].table.cell(k + 1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE              
            tables_i[0].table.cell(k + 1, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            lsl, usl = '{:.4f}'.format(meas_list[i * num_rows_meas + k].loc[0,'LSpecL']), '{:.4f}'.format(meas_list[i * num_rows_meas + k].loc[0,'USpecL'])
            tables_i[0].table.cell(k + 1, 2).text = str(lsl) + '/' + str(usl) 
            tables_i[0].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[0].table.cell(k + 1, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
            fallout = fallouts[i * num_rows_meas + k]
            if str(lsl) == 'nan':
                tables_i[0].table.cell(k + 1, 3).text = 'N/A'
            else:
                # if lsl us not empty, filling fallout cells with color
                tables_i[0].table.cell(k + 1, 3).text = '{:.1%}'.format(fallout)   
                fill = tables_i[0].table.cell(k + 1, 3).fill
                fore_color = fill.solid()
                if fallout == 0:
                    tables_i[0].table.cell(k + 1, 3).text = 'Pass'
                    fill.fore_color.rgb = RGBColor(198, 239, 206)
                elif fallout < 0.05:
                    fill.fore_color.rgb = RGBColor(255, 235, 156)
                elif fallout >= 0.05:
                    fill.fore_color.rgb = RGBColor(255, 199, 206)
            
            tables_i[0].table.cell(k+1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER    
            tables_i[0].table.cell(k+1, 3).vertical_anchor = MSO_ANCHOR.MIDDLE
            tables_i[0].table.cell(k + 1, 4).text = '{:.4f}'.format(ranges[i * num_rows_meas + k])
            tables_i[0].table.cell(k+1, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tables_i[0].table.cell(k + 1, 4).vertical_anchor = MSO_ANCHOR.MIDDLE
                                    
    table_style_12.write_shape(tables_i[0])
    
# slides with measurements under 1 description and worst measurements
dop_sl = 0 # 1 description - 2 slides plus several more if list of measurements doesn't fit to 1 slide
# dop_sl is the sum of added slides to output all the measurements in list
for i in range(num_req):
    num_meas_i = len(meas_list[i]) # how many measurements in certain description
    if num_meas_i <= num_rows_big_tabl: # condition if they can ba placed at 1 slide or we need several
        copied_slide = duplicate_slide(pres, 3)
        slide_loc = num_templates + 2 * init_slides_num + num_idx + 2*i + dop_sl # right location is highly important
        #print(i, slide_loc) # so better check slides allocation
        slide_i = pres.slides[slide_loc] 
        tables_i = [shape for shape in slide_i.shapes if shape.has_table]
        
        # header 
        if with_index:
            tables_i[0].table.cell(0, 1).text = meas_list[i].loc[0,'Description_y'].split(':')[0]+' -'+meas_list[i].loc[0,'Description_y'].split(':')[1]
        else:
            tables_i[0].table.cell(0, 1).text = str(meas_list[i].loc[0,'Description_y'])
        
        # yellow text frames - meas names
        title_shapes = [shape for shape in slide_i.shapes if shape.has_text_frame]
        titles = [shape for shape in title_shapes if shape.has_text_frame]

        for k in range(len(meas_list[i])):
            tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
            text_01 = meas_list[i].loc[k, 'Name']
            titles[k + 2].text = meas_list[i].loc[k, 'Name'] # starting from index of 2 as there's other text frames on slide
            # font formatting of text boxes
            font_yllw_box.write_paragraph(titles[k + 2].text_frame.paragraphs[0])
            tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Nominal'])
            tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
            tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
            tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
            tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
            tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
            tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
            tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
            tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
            # alignment
            for w in range(10):
                tables_i[1].table.cell(k+1, w).vertical_anchor = MSO_ANCHOR.MIDDLE                
                     
        table_style_12.write_shape(tables_i[0])
        table_style_10.write_shape(tables_i[1])
        paragraph_style.write_shape(tables_i[1])
        # meas name alignment on the left side
        for k in range(len(meas_list[i])):
            tables_i[1].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    else: # we need several slides per 1 description
        if num_meas_i % num_rows_big_tabl == 0:
            slides_per_req = int(num_meas_i / num_rows_big_tabl) # how many such slides we need
        else:
            slides_per_req = num_meas_i // num_rows_big_tabl + 1
            
        for j in range(slides_per_req):
            copied_slide = duplicate_slide(pres, 3)  
            slide_loc = num_templates + 2 * init_slides_num + num_idx + 2*i + dop_sl + j
            slide_i = pres.slides[slide_loc] 
            #print(i, slide_loc)
            #print(len(pres.slides))
            tables_i = [shape for shape in slide_i.shapes if shape.has_table]
            title_shapes = [shape for shape in slide_i.shapes if shape.has_text_frame]
            titles = [shape for shape in title_shapes if shape.has_text_frame]
        
            if with_index:
                tables_i[0].table.cell(0, 1).text = meas_list[i].loc[0,'Description_y'].split(':')[0]+' -'+meas_list[i].loc[0,'Description_y'].split(':')[1]
            else:
                tables_i[0].table.cell(0, 1).text = str(meas_list[i].loc[0,'Description_y'])
            
            # we just divide df to smaller ones to be put on slides' tables one by one
            if j != slides_per_req - 1:  
                meas_list_j = meas_list[i].iloc[num_rows_big_tabl * j: num_rows_big_tabl * (j+1), ]
            else:
                meas_list_j = meas_list[i].iloc[num_rows_big_tabl * j: , ]
            meas_list_j.reset_index(inplace=True, drop=True) # we consider them as new smaller df-tables
                
            for k in range(len(meas_list_j)):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                titles[k + 2].text = meas_list_j.loc[k, 'Name']
                font_yllw_box.write_paragraph(titles[k + 2].text_frame.paragraphs[0])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Nominal'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])

                table_style_12.write_shape(tables_i[0])
                table_style_10.write_shape(tables_i[1])
        
            paragraph_style.write_shape(tables_i[1])
            for k in range(len(meas_list_j)):
                tables_i[1].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        # as we firstly assume: 1 description = 1 meas list slide, so we count only added ones       
        dop_sl += slides_per_req - 1  
          
    # indexation is the same for meas_list and contr_choice, so following order of descriptions
    # taking worst measurement (contr) for meas_list[i]
    contr = contr_choice[i]
    loc = 0 
    # where it is located in contr_names_list
    # (array of all! measure-contributors tables, order - how they are placed in rss)
    for m in range(len(contr_names_list)):
        if contr_names_list[m] == contr:
            loc = m
    subtabl = sublist_dfs[loc] # we've chosen df for worst measurement
    
    copied_slide = duplicate_slide(pres, 4) # template for contributors slide
    slide_loc = num_templates + 2 * init_slides_num + num_idx + 2*i + 1 + dop_sl
    slide_i = pres.slides[slide_loc] 
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]

    lsl, usl = '{:.4f}'.format(meas_list[i].loc[0,'LSpecL']), '{:.4f}'.format(meas_list[i].loc[0,'USpecL'])
    
    # blue header
    if with_index:
        tables_i[4].table.cell(0,0).text = meas_list[i].loc[0,'Description_y'].split(sep=':')[0] + ' - '+ contr + meas_list[i].loc[0,'Description_y'].split(sep=':')[1] + '. Spec Limits: ' + lsl+'/'+usl
    else:
        tables_i[4].table.cell(0,0).text = contr+ ' - '+meas_list[i].loc[0,'Description_y'] + '. Spec Limits: ' + lsl+'/'+usl
    style_contr_header.write_shape(tables_i[4])
    
    title_shapes = [shape for shape in slide_i.shapes if shape.has_text_frame]
    titles = [shape for shape in title_shapes if shape.has_text_frame]
    
    # filling yellow text boxes    
    titles[3].text = lsl
    titles[2].text = usl
    titles[4].text = 'Nom'
    
    # filling table with major contributors
    # placing number of contr-s <= rows in table
    len_tabl = len(subtabl) if len(subtabl) < num_rows_contr else num_rows_contr 
    for k in range(len_tabl):
        tables_i[0].table.cell(k+1, 0).text = str(subtabl.loc[k, 'contr_name'])
        tables_i[0].table.cell(k+1, 1).text = '{:.4f}'.format(float(subtabl.loc[k, 'tol_range']))
        tables_i[0].table.cell(k+1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        proc = '{}% '.format(subtabl.loc[k, 'contr'][:-4])
        tables_i[0].table.cell(k+1, 2).text = proc
        titles[5+k].text = proc # put procent of contribution to text box
        tables_i[0].table.cell(k+1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table_style_10.write_shape(tables_i[0])
    
    idx = contr_locs[i] # where worst measurement is located in df of its description
    tables_i[1].table.cell(1, 0).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Nominal'])
    tables_i[1].table.cell(1, 1).text = '{:.1%}'.format(meas_list[i].loc[idx,'L-OUT'])
    tables_i[1].table.cell(1, 2).text = '{:.1%}'.format(meas_list[i].loc[idx,'H-OUT'])
    tables_i[1].table.cell(1, 3).text = '{:.1%}'.format(meas_list[i].loc[idx, 'Tot-OUT'])
    tables_i[2].table.cell(1, 0).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Est.Low'])
    tables_i[2].table.cell(1, 1).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Est.High'])
    tables_i[3].table.cell(1, 0).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Est.Low.C'])
    tables_i[3].table.cell(1, 1).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Est.High.C'])
    tables_i[2].table.cell(1, 2).text = '{:.4f}'.format(meas_list[i].loc[idx, 'Est.Range'])
    table_style_10.write_shape(tables_i[1])
    table_style_10.write_shape(tables_i[2])
    table_style_10.write_shape(tables_i[3])
    paragraph_style.write_shape(tables_i[1])
    paragraph_style.write_shape(tables_i[2])
    paragraph_style.write_shape(tables_i[3])
    
    for filename in os.listdir('images/'):
        if contr in filename:
            img_path = os.path.join('images/', filename)        
    img_i = slide_i.shapes.add_picture(img_path, left, top, width)
     
# all slides we've generated so far
num_slides_main =  num_templates + 2 * init_slides_num + 2 * num_req + num_idx + dop_sl

# number of slides needed for data from csv
num_end_slides = int(num_meas_total / num_rows_end) if num_meas_total % num_rows_end == 0 else num_meas_total // num_rows_end + 1

# tables from csv put in the end
for i in range(num_end_slides):
    copied_slide = duplicate_slide(pres, num_templates-1)  
    slide_i = pres.slides[num_slides_main + i]
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]  
    
    # we divide our csv-df to small dfs to be put on slides and fill table with them
    if i != num_end_slides - 1:  
        list_i = df1_init.iloc[num_rows_end * i: num_rows_end * (i+1), ] # last dataframe
    else:
        list_i = df1_init.iloc[num_rows_end * i: , ]
    list_i.reset_index(inplace=True, drop=True) # important 
    
    for k in range(len(list_i)):
        tables_i[0].table.cell(k+1, 0).text = str(list_i.loc[k, 'Name'])
        tables_i[0].table.cell(k+1, 1).text = '{:.4f}'.format(list_i.loc[k, 'Nominal'])
        tables_i[0].table.cell(k+1, 2).text = '{:.4f}'.format(list_i.loc[k, 'Mean'])
        tables_i[0].table.cell(k+1, 3).text = '{:.4f}'.format(list_i.loc[k, '6-Sigma'])
        tables_i[0].table.cell(k+1, 4).text = '{:.4f}'.format(list_i.loc[k, 'Min'])
        tables_i[0].table.cell(k+1, 5).text = '{:.4f}'.format(list_i.loc[k, 'Max'])
        tables_i[0].table.cell(k+1, 6).text = '{:.4f}'.format(list_i.loc[k, 'Range'])
        lsl, usl = list_i.loc[k, 'LSL'], list_i.loc[k, 'USL']
        if lsl == lsl_nan:
            tables_i[0].table.cell(k+1, 7).text = 'nan'
        else:
            tables_i[0].table.cell(k+1, 7).text = '{:.4f}'.format(lsl) 
        if usl == usl_nan:    
            tables_i[0].table.cell(k+1, 8).text = 'nan'
        else:
            tables_i[0].table.cell(k+1, 8).text = '{:.4f}'.format(usl)
        tables_i[0].table.cell(k+1, 9).text = '{:.1%}'.format(list_i.loc[k,'L-OUT'])
        tables_i[0].table.cell(k+1, 10).text = '{:.1%}'.format(list_i.loc[k,'H-OUT'])
        tables_i[0].table.cell(k+1, 11).text = '{:.1%}'.format(list_i.loc[k, 'Tot-OUT'])
        etype = list_i.loc[k, 'Est.Type']
        # somehow there's white spaces we don't need
        if etype == '   Normal':
            tables_i[0].table.cell(k+1, 12).text = etype[3:]
        else:
            tables_i[0].table.cell(k+1, 12).text = etype
        tables_i[0].table.cell(k+1, 13).text = '{:.4f}'.format(list_i.loc[k, 'Est.Low'])
        tables_i[0].table.cell(k+1, 14).text = '{:.4f}'.format(list_i.loc[k, 'Est.High'])
        tables_i[0].table.cell(k+1, 15).text = '{:.4f}'.format(list_i.loc[k, 'Est.Range'])
        
    paragraph_style.write_shape(tables_i[0])
    # alignment of meas name
    for k in range(len(list_i)):
                tables_i[0].table.cell(k+1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT        
    table_style_10.write_shape(tables_i[0])      
 
# deleting template slides from the beginning of the pres, no changes here please       
for i in range(num_templates-1,-1,-1) : 
    rId = pres.slides._sldIdLst[i].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[i]

pres.save('pres_latch.pptx') # saving file  