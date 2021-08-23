#!/usr/bin/env python
# coding: utf-8
num_rows_0 = 15
num_rows_2 = 12
num_rows_3 = 4

# which initial template is used for test_pres
initial_template = 1

tabl3_if_string = 1 # if in cells = 0 

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.creator import PPTXCreator
from pptx_tools.position import PPTXPosition
from pptx_tools.table_style import PPTXTableStyle
import pandas as pd
import six
import copy
import numpy as np
import os

# formatting font in tables
table_style = PPTXTableStyle()
table_style.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=12)

table_style1 = PPTXTableStyle()
table_style1.font_style = PPTXFontStyle().set(italic=False, name="Arial", size=10)

def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[5]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )
    return copied_slide

# data to be put in table (x-measurements are dropped)
df0 = pd.read_excel('Book1.xlsx', sheet_name='Measures', header=0)
df0 = df0.drop_duplicates(subset=['PART', 'Measure'])
df0.reset_index(inplace=True, drop=True)

wrong_ind = []
for i in range(len(df0)):
    if 'x' in df0.iloc[i,1]:
        wrong_ind.append(i)
    if df0.loc[i, 'As Output'] == 0:
        wrong_ind.append(i)
    if df0.loc[i, 'Active'] == 0:
        wrong_ind.append(i)
df0.drop(wrong_ind, inplace=True)
df0.reset_index(inplace=True, drop=True)

df0 = df0[['Description', 'Type', 'Measure', 'PART', 'USpecL', 'LSpecL', 'Relative/Abs']]

if df0.loc[0, 'Description'][0] == '#':
    with_index = True
else:
    with_index = False
    
if with_index:
    for i in range(len(df0)):
        desc = str(df0.loc[i, 'Description']).split(sep=':')
        df0.loc[i,'idx'] = int(desc[0][1:])

df1 = pd.read_excel('tabl2.xlsx', header=12)

wrong_ind = []
for i in range(len(df1)):
    if 'x' in df1.iloc[i,1]:
        wrong_ind.append(i)
df1.drop(wrong_ind, inplace=True)
df1.reset_index(inplace=True, drop=True)
df1 = df1[['Name', 'Description', 'Nominal', 'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range']]

# УБРАТЬ РЕШЕТКИ В ДВУХ СТРОКАХ НИЖЕ, ЕСЛИ ИСПОЛЬЗУЕТСЯ ALIAS В ИЗМЕРЕНИЯХ
#for i in range(len(df1)):
#    df1.loc[i, 'Name'] = df1.loc[i, 'Name'][3:]

# merging data about certain measurements
df_01 = df0.merge(df1, left_on='Measure', right_on='Name')
df_01 = df_01.drop(['Description_x', 'Measure'], axis=1)

if with_index:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs', 'idx']]
    df_01['Est.Low.C'] = df_01['Est.Low'] - df_01['Nominal']
    df_01['Est.High.C'] = df_01['Est.High'] - df_01['Nominal']
    df_01 = df_01.sort_values(by='idx', ignore_index=True)
else:
    df_01 = df_01[['Name', 'PART', 'Description_y', 'Type', 'USpecL', 'LSpecL', 'Nominal', 
        'L-OUT', 'H-OUT', 'Tot-OUT', 'Est.Low', 'Est.High', 'Est.Range', 'Relative/Abs']]
    df_01['Est.Low.C'] = df_01['Est.Low'] - df_01['Nominal']
    df_01['Est.High.C'] = df_01['Est.High'] - df_01['Nominal']

df2 = pd.read_excel('tabl3.xlsx', skiprows=10)


# Вторая табличка:

# array of tables corresponding to measurements
# and array of descriptions for convenience

meas_list = []
meas_names = []

for name in df_01['Description_y'].unique():
    meas_names.append(name)
    df_i = df_01[df_01['Description_y'] == name]
    df_i.reset_index(inplace=True, drop=True)
    meas_list.append(df_i)


# creating and filling template_choice
temp_array = [initial_template for i in range(len(meas_names))]
links = pd.DataFrame({'req':meas_names, 'template':temp_array})
links.to_excel('./template_choice.xlsx')


# Решаем, какие точки класть в третью таблицу:

contr_choice = []
fallouts = []

for df in meas_list:
    if len(df) == 1:
        point = df.loc[0, 'Name']
    else:
        ind = 0
        max_t_o = df['Tot-OUT'].max()
        for val in df['Tot-OUT']:
            if val == max_t_o:
                ind += 1
        if ind > 1:
            point = df.loc[df['Est.Range'].idxmax(), 'Name']
        else:
            point = df.loc[df['Tot-OUT'].idxmax(), 'Name']
            
    contr_choice.append(point)
    fallouts.append(df['Tot-OUT'].max())

# working with last table
# finding empty lines 
# getting array of indexes where records start

newrec_arr = []
for i in range(len(df2)):
    if str(df2.iloc[i, 0]) == 'nan':
        newrec_arr.append(i)

l = newrec_arr
l_mod = [0] + l + [len(df2) + 1]  

# info on every measurement
list_dfs = [df2.iloc[l_mod[n]:l_mod[n+1]] for n in range(len(l_mod)-1)]

for i in range(1, len(list_dfs)):
    list_dfs[i].reset_index(inplace=True, drop=True)
    list_dfs[i].drop([0], inplace=True)
    list_dfs[i].reset_index(inplace=True, drop=True)


contr_names_list = []
sublist_dfs = []
for df in list_dfs:
    # в случае ALIAS используем следующую строку, комментируем через одну
    # contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][3:-1])
    contr_names_list.append(df.iloc[0, 0].split(sep=' ')[1][:-1])
    num_recs = int((len(df) - 5) / 2)
    
    sub_df = pd.DataFrame(index=range(num_recs), columns=['contr_name', 'type', 'tol_range', 'contr'])
    for i in range(num_recs):
         if tabl3_if_string == 1:
             str_i = df.iloc[4 + i * 2, 0].split()
             
             sub_df.loc[i, 'tol_range'] = str_i[4][2:] # range
             sub_df.loc[i, 'contr'] = str_i[5] # contribution
            
             typ = str_i[1][:2]
             if typ == 'CR':
                 sub_df.loc[i, 'type'] = 'Positional'
             elif typ == 'CH' or typ == 'CP':
                 sub_df.loc[i, 'type'] = 'Size'
             else:
                 sub_df.loc[i, 'type'] = 'Linear'
    
             if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan':
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
             else:
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]
                 
         else:
             str_i = df.iloc[4 + i * 2, 1]
             sub_df.loc[i, 'tol_range'] = df.iloc[4 + i * 2, 4][2:] # range
             sub_df.loc[i, 'contr'] = str(df.iloc[4 + i * 2, 5]) # contribution
             
             typ = str_i[:1]
             if typ == 'CR':
                 sub_df.loc[i, 'type'] = 'Positional'
             elif typ == 'CH' or typ == 'CP':
                 sub_df.loc[i, 'type'] = 'Size'
             else:
                 sub_df.loc[i, 'type'] = 'Linear'
     
             if str(df.iloc[4 + i * 2 + 1, 1]) != 'nan':
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]+ ',' + df.iloc[4 + i * 2 + 1, 1]
             else:
                 sub_df.loc[i, 'contr_name'] = df.iloc[4 + i * 2 + 1, 0][3:]

    sublist_dfs.append(sub_df)   

# opening existing presentation (template with tables)
# генерация заданного количества слайдов

pres = Presentation('boe_template.pptx')
pres.slide_width = Inches(13.333333) # widescreen
pres.slide_height = Inches(7.5)

slides = [slide for slide in pres.slides] # array of slide-objects


num_templates = 10
num_req = len(meas_list) # number of req = number of slides
# array of template numbers corresponding to list of descriptions meas_names
links = pd.read_excel('template_choice.xlsx', header=0)

if num_req % num_rows_0 == 0:
    init_slides_num = int(num_req / num_rows_0)
else:
    init_slides_num = num_req // num_rows_0 + 1

for i in range(init_slides_num):
    copied_slide = duplicate_slide(pres, 0)  
    slide_i = pres.slides[num_templates + i]
    tables_i = [shape for shape in slide_i.shapes if shape.has_table]
    for k in range(num_rows_0):
        if k <= num_req - num_rows_0 * i - 1:
            tables_i[0].table.cell(k + 1, 1).text = str(meas_names[i * num_rows_0 + k])
            lsl, usl = '{:.4f}'.format(meas_list[i * num_rows_0 + k].iloc[0,5]), '{:.4f}'.format(meas_list[i * num_rows_0 + k].iloc[0,4])
            tables_i[0].table.cell(k + 1, 2).text = str(lsl) + '/' + str(usl) 
            fallout = fallouts[i * num_rows_0 + k]
            
            if str(lsl) == 'nan':
                tables_i[0].table.cell(k + 1, 3).text = 'nan'
            else:
                tables_i[0].table.cell(k + 1, 3).text = '{:.1%}'.format(fallout)   
                fill = tables_i[0].table.cell(k + 1, 3).fill
                fore_color = fill.solid()
                if fallout == 0:
                    tables_i[0].table.cell(k + 1, 3).text = 'Pass'
                    fill.fore_color.rgb = RGBColor(198, 239, 206)
                elif fallout < 0.05:
                    fill.fore_color.rgb = RGBColor(255, 235, 156)
                elif fallout >= 0.05:
                    fill.fore_color.rgb = RGBColor(255, 199, 206)
                                    
    table_style.write_shape(tables_i[0])


# filling tables slide by slide
dop_sl = 0

for i in range(num_req):
    num_meas_i = len(meas_list[i])
    if num_meas_i <= num_rows_2:
        template_type = links['template'][i]
        copied_slide = duplicate_slide(pres, template_type)
        slide_loc = num_templates + init_slides_num + i + dop_sl
        #print(i, slide_loc)
        slide_i = pres.slides[slide_loc] 
        tables_i = [shape for shape in slide_i.shapes if shape.has_table]
        
        req, mes_type = meas_list[i].iloc[0,2], meas_list[i].iloc[0, 3]
        tables_i[0].table.cell(0,2).text = str(req)
        tables_i[0].table.cell(1,2).text = str(mes_type)
        if meas_list[i].loc[0, 'Relative/Abs'] == 0:
            tables_i[0].table.cell(1,8).text = 'Absolute' 
        if meas_list[i].loc[0, 'Relative/Abs'] == 1:
            tables_i[0].table.cell(1,8).text = 'Relative'
        lsl, usl = round(meas_list[i].iloc[0,5], 4), round(meas_list[i].iloc[0,4], 4)
        tables_i[0].table.cell(1,6).text = str(lsl) + '/' + str(usl)   

        if links['template'][i] == 1:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = str(round(meas_list[i].loc[k, 'Nominal'], 4))
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 2:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 3:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 4:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 5:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 6:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])       
        elif links['template'][i] == 7:
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = str(round(meas_list[i].loc[k, 'Nominal'], 4))
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
        elif links['template'][i] == 8:   
            for k in range(len(meas_list[i])):
                tables_i[1].table.cell(k+1, 0).text = str(meas_list[i].loc[k, 'Name'])
                tables_i[1].table.cell(k+1, 1).text = str(round(meas_list[i].loc[k, 'Nominal'], 4))
                tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list[i].loc[k,'L-OUT'])
                tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list[i].loc[k,'H-OUT'])
                tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list[i].loc[k, 'Tot-OUT'])
                tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Low.C'])
                tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.High.C'])
                tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list[i].loc[k, 'Est.Range'])
    
        contr = contr_choice[i]
        loc = 0
        for m in range(len(contr_names_list)):
            if contr_names_list[m] == contr:
                loc = m
        subtabl = sublist_dfs[loc] 
        tables_i[2].table.cell(0, 0).text = 'Major Contributors to Variations ({})'.format(contr)
        len_tabl = len(subtabl) if len(subtabl) < num_rows_3 else num_rows_3
        for k in range(len_tabl):
            tables_i[2].table.cell(k+1, 0).text = str(subtabl.iloc[k, 0])
            tables_i[2].table.cell(k+1, 1).text = str(subtabl.iloc[k, 1])
            tables_i[2].table.cell(k+1, 2).text = str(subtabl.iloc[k, 2])
            if tabl3_if_string == 0:
                tables_i[2].table.cell(k+1, 3).text = '{:.1%} '.format(float(subtabl.iloc[k, 3]))
            else:
                tables_i[2].table.cell(k+1, 3).text = '{}% '.format(subtabl.iloc[k, 3][:-4])

        table_style1.write_shape(tables_i[0])
        table_style1.write_shape(tables_i[1])
        table_style1.write_shape(tables_i[2])

        for filename in os.listdir('images/'):
            if contr in filename:
                img_path = os.path.join('images/', filename)        
        left = Inches(9.5)
        top = Inches(5)
        width = Inches(3)
        img_i = slide_i.shapes.add_picture(img_path, left, top, width)
        
    else:
        if num_meas_i % num_rows_2 == 0:
            slides_per_req = int(num_meas_i / num_rows_2)
        else:
            slides_per_req = num_meas_i // num_rows_2 + 1
            
        for j in range(slides_per_req):
            template_type = links['template'][i]
            copied_slide = duplicate_slide(pres, template_type)  
            slide_loc = num_templates + init_slides_num + i + dop_sl + j
            slide_i = pres.slides[slide_loc] 
            #print(i, slide_loc)
            tables_i = [shape for shape in slide_i.shapes if shape.has_table]
            
            req, mes_type = meas_list[i].iloc[0,2], meas_list[i].iloc[0, 3]
            tables_i[0].table.cell(0,2).text = str(req)
            tables_i[0].table.cell(1,2).text = str(mes_type)
            if meas_list[i].loc[0, 'Relative/Abs'] == 0:
                tables_i[0].table.cell(1,8).text = 'Absolute' 
            if meas_list[i].loc[0, 'Relative/Abs'] == 1:
                tables_i[0].table.cell(1,8).text = 'Relative'
            lsl, usl = round(meas_list[i].iloc[0,5], 4), round(meas_list[i].iloc[0,4], 4)
            tables_i[0].table.cell(1,6).text = str(lsl) + '/' + str(usl) 
            
            contr = contr_choice[i]
            loc = 0
            for m in range(len(contr_names_list)):
                if contr_names_list[m] == contr:
                    loc = m
            subtabl = sublist_dfs[loc] 
            tables_i[2].table.cell(0, 0).text = 'Major Contributors to Variations ({})'.format(contr)
            len_tabl = len(subtabl) if len(subtabl) < num_rows_3 else num_rows_3
            for k in range(len_tabl):
                tables_i[2].table.cell(k+1, 0).text = str(subtabl.iloc[k, 0])
                tables_i[2].table.cell(k+1, 1).text = str(subtabl.iloc[k, 1])
                tables_i[2].table.cell(k+1, 2).text = str(subtabl.iloc[k, 2])
                if tabl3_if_string == 0:
                    tables_i[2].table.cell(k+1, 3).text = '{:.1%} '.format(float(subtabl.iloc[k, 3]))
                else:
                    tables_i[2].table.cell(k+1, 3).text = '{}% '.format(subtabl.iloc[k, 3][:-4])

            for filename in os.listdir('images/'):
                if contr in filename:
                    img_path = os.path.join('images/', filename)        
            left = Inches(9.5)
            top = Inches(5)
            width = Inches(3)
            img_i = slide_i.shapes.add_picture(img_path, left, top, width)
            
            if j != slides_per_req - 1:  
                meas_list_j = meas_list[i].iloc[num_rows_2 * j: num_rows_2 * (j+1), ]
            else:
                meas_list_j = meas_list[i].iloc[num_rows_2 * j: , ]
            meas_list_j.reset_index(inplace=True, drop=True)
            
            if links['template'][i] == 1:
                
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.iloc[k, 0])
                    tables_i[1].table.cell(k+1, 1).text = str(round(meas_list_j.iloc[k, 6], 4))
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.iloc[k, 7])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.iloc[k, 8])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.iloc[k, 9])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.iloc[k, 10])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.iloc[k, 11])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.iloc[k, 14])
                    tables_i[1].table.cell(k+1, 8).text = '{:.4f}'.format(meas_list_j.iloc[k, 15])
                    tables_i[1].table.cell(k+1, 9).text = '{:.4f}'.format(meas_list_j.iloc[k, 12])
            if links['template'][i] == 2:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 3:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 4:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 5:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 6:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])       
            elif links['template'][i] == 7:
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = str(round(meas_list_j.loc[k, 'Nominal'], 4))
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High'])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])
            elif links['template'][i] == 8:   
                for k in range(len(meas_list_j)):
                    tables_i[1].table.cell(k+1, 0).text = str(meas_list_j.loc[k, 'Name'])
                    tables_i[1].table.cell(k+1, 1).text = str(round(meas_list_j.loc[k, 'Nominal'], 4))
                    tables_i[1].table.cell(k+1, 2).text = '{:.1%}'.format(meas_list_j.loc[k,'L-OUT'])
                    tables_i[1].table.cell(k+1, 3).text = '{:.1%}'.format(meas_list_j.loc[k,'H-OUT'])
                    tables_i[1].table.cell(k+1, 4).text = '{:.1%}'.format(meas_list_j.loc[k, 'Tot-OUT'])
                    tables_i[1].table.cell(k+1, 5).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Low.C'])
                    tables_i[1].table.cell(k+1, 6).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.High.C'])
                    tables_i[1].table.cell(k+1, 7).text = '{:.4f}'.format(meas_list_j.loc[k, 'Est.Range'])        
            
            table_style1.write_shape(tables_i[0])
            table_style1.write_shape(tables_i[1])
            table_style1.write_shape(tables_i[2])
        dop_sl += slides_per_req - 1    
            
    
# deleting several template slides in the beginning
for i in range(num_templates-1,-1,-1) : 
    rId = pres.slides._sldIdLst[i].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[i]


pres.save('test_pres.pptx') # saving file   




